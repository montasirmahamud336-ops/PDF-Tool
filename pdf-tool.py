from __future__ import annotations

"""PDFForge FastAPI application for PDF, image, and document workflows."""

import base64
import os
import re
import sys
import tempfile
import uuid
import webbrowser
from collections import Counter
from contextlib import redirect_stderr, redirect_stdout
from datetime import datetime
from html import escape
from io import BytesIO, StringIO
from pathlib import Path
from typing import List, Optional

import uvicorn
from fastapi import BackgroundTasks, FastAPI, File, Form, HTTPException, UploadFile
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse
from PIL import Image
from PyPDF2 import PdfMerger, PdfReader, PdfWriter
from pdf2image import convert_from_bytes
from pdf2image.exceptions import PDFInfoNotInstalledError, PDFPageCountError, PDFSyntaxError

DEPS_DIR = Path(__file__).with_name(".deps")
if DEPS_DIR.exists():
    sys.path.insert(0, str(DEPS_DIR))

try:
    import pypdfium2 as pdfium
except Exception:
    pdfium = None
try:
    from reportlab.pdfgen import canvas as rl_canvas
    from reportlab.lib.pagesizes import letter
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer
except Exception:
    rl_canvas = None
    letter = None
    getSampleStyleSheet = None
    Paragraph = None
    SimpleDocTemplate = None
    Spacer = None
try:
    from docx import Document
except Exception:
    Document = None
try:
    from pptx import Presentation
except Exception:
    Presentation = None
try:
    from openpyxl import Workbook, load_workbook
except Exception:
    Workbook = None
    load_workbook = None
try:
    with StringIO() as _weasyprint_output, redirect_stdout(_weasyprint_output), redirect_stderr(_weasyprint_output):
        from weasyprint import HTML
except Exception:
    HTML = None
try:
    import fitz  # PyMuPDF
except Exception:
    fitz = None
try:
    from googletrans import Translator
except Exception:
    Translator = None


APP_DIR = Path(__file__).resolve().parent
APP_TITLE = "PDFForge Tool API"
APP_DESCRIPTION = "Multi-tool PDF workflows for merge, split, convert, scan, and document utilities."
APP_VERSION = "1.0.0"
SUPPORTED_MERGE_EXTS = {
    ".pdf",
    ".png",
    ".jpg",
    ".jpeg",
    ".webp",
    ".html",
    ".htm",
    ".docx",
    ".pptx",
    ".xlsx",
    ".txt",
    ".md",
    ".csv",
    ".json",
}
SUPPORTED_IMAGE_FORMATS = {"png", "jpeg", "jpg", "webp"}
TEXT_SOURCE_EXTS = {".txt", ".md", ".csv", ".json"}

def _env_flag(name: str, default: bool = False) -> bool:
    raw = os.getenv(name)
    if raw is None:
        return default
    return raw.strip().lower() in {"1", "true", "yes", "on"}


def _is_hosted_environment() -> bool:
    hosted_markers = ("VERCEL", "RAILWAY_ENVIRONMENT", "RENDER", "K_SERVICE", "DYNO")
    return any(os.getenv(marker) for marker in hosted_markers)


def _is_writable_dir(path: Path) -> bool:
    try:
        path.mkdir(parents=True, exist_ok=True)
        probe = path / ".write-test"
        probe.write_text("ok", encoding="utf-8")
        probe.unlink(missing_ok=True)
        return True
    except Exception:
        return False


def _resolve_runtime_root() -> Path:
    override = os.getenv("PDF_TOOL_RUNTIME_DIR")
    candidates: List[Path] = []
    if override:
        candidates.append(Path(override).expanduser())

    temp_root = Path(tempfile.gettempdir()) / "pdf-tool-runtime"
    if os.getenv("VERCEL"):
        candidates.append(temp_root)
        candidates.append(APP_DIR / ".runtime")
    else:
        candidates.append(APP_DIR)
        candidates.append(APP_DIR / ".runtime")
        candidates.append(temp_root)

    for candidate in candidates:
        if _is_writable_dir(candidate):
            return candidate

    raise RuntimeError("Could not find a writable runtime directory for uploads and outputs.")


def _parse_cors_origins() -> List[str]:
    raw = os.getenv("PDF_TOOL_CORS_ORIGINS", "*").strip()
    if not raw or raw == "*":
        return ["*"]
    return [origin.strip() for origin in raw.split(",") if origin.strip()]


RUNTIME_ROOT = _resolve_runtime_root()
UPLOAD_DIR = RUNTIME_ROOT / "uploads"
OUTPUT_DIR = RUNTIME_ROOT / "outputs"
INDEX_FILE = APP_DIR / "index.html"

UPLOAD_DIR.mkdir(parents=True, exist_ok=True)
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

app = FastAPI(
    title=APP_TITLE,
    description=APP_DESCRIPTION,
    version=APP_VERSION,
)
app.add_middleware(
    CORSMiddleware,
    allow_origins=_parse_cors_origins(),
    allow_methods=["*"],
    allow_headers=["*"],
)


def _safe_remove(path: str | Path) -> None:
    try:
        Path(path).unlink(missing_ok=True)
    except Exception:
        pass


def _save_upload(upload: UploadFile, target_dir: Path) -> Path:
    ext = Path(upload.filename or "").suffix.lower()
    unique = f"{uuid.uuid4().hex}{ext}"
    path = target_dir / unique
    with path.open("wb") as f:
        f.write(upload.file.read())
    return path


def _image_to_pdf(image_path: Path) -> Path:
    output_path = OUTPUT_DIR / f"img_as_pdf_{uuid.uuid4().hex}.pdf"
    with Image.open(image_path) as img:
        if img.mode in ("RGBA", "LA"):
            bg = Image.new("RGB", img.size, "white")
            bg.paste(img, mask=img.getchannel("A"))
            ready = bg
        elif img.mode == "P" and "transparency" in img.info:
            rgba = img.convert("RGBA")
            bg = Image.new("RGB", rgba.size, "white")
            bg.paste(rgba, mask=rgba.getchannel("A"))
            ready = bg
        else:
            ready = img.convert("RGB")
        ready.save(output_path, "PDF")
    return output_path


def _read_text_source(path: Path) -> str:
    for encoding in ("utf-8", "utf-8-sig", "latin-1"):
        try:
            return path.read_text(encoding=encoding)
        except UnicodeDecodeError:
            continue
    return path.read_text(encoding="utf-8", errors="ignore")


def _split_text_blocks(text: str) -> List[str]:
    blocks = [block.strip() for block in re.split(r"\n\s*\n", text or "") if block.strip()]
    if blocks:
        return blocks
    fallback = [line.strip() for line in (text or "").splitlines() if line.strip()]
    return fallback or ["No readable content was found."]


def _pdf_split_to_files(pdf_bytes: bytes, page: int) -> tuple[Path, Path, int]:
    reader = PdfReader(BytesIO(pdf_bytes))
    total_pages = len(reader.pages)
    if total_pages < 2:
        raise HTTPException(status_code=400, detail="PDF must have at least 2 pages to split.")
    if page < 1 or page >= total_pages:
        raise HTTPException(
            status_code=400,
            detail=f"Page number must be between 1 and {total_pages - 1}.",
        )

    part1 = PdfWriter()
    part2 = PdfWriter()

    for i in range(page):
        part1.add_page(reader.pages[i])
    for i in range(page, total_pages):
        part2.add_page(reader.pages[i])

    part1_path = OUTPUT_DIR / f"part1_{uuid.uuid4().hex}.pdf"
    part2_path = OUTPUT_DIR / f"part2_{uuid.uuid4().hex}.pdf"

    with part1_path.open("wb") as f:
        part1.write(f)
    with part2_path.open("wb") as f:
        part2.write(f)

    return part1_path, part2_path, total_pages


def _resolve_image_format(fmt: str) -> tuple[str, str]:
    normalized = (fmt or "png").strip().lower()
    if normalized not in SUPPORTED_IMAGE_FORMATS:
        raise HTTPException(
            status_code=400,
            detail="Unsupported format. Use png, jpeg, jpg, or webp.",
        )

    if normalized in {"jpeg", "jpg"}:
        return normalized, "JPEG"
    if normalized == "webp":
        return normalized, "WEBP"
    return normalized, "PNG"


def _read_pdf(pdf_bytes: bytes, password: str | None = None, strict: bool = True) -> PdfReader:
    if not pdf_bytes:
        raise HTTPException(status_code=400, detail="Uploaded PDF is empty.")
    try:
        reader = PdfReader(BytesIO(pdf_bytes), strict=strict)
    except Exception as exc:
        raise HTTPException(status_code=400, detail="Invalid PDF file.") from exc

    if reader.is_encrypted:
        if not password:
            raise HTTPException(status_code=400, detail="PDF is password-protected.")
        try:
            ok = reader.decrypt(password)
        except Exception:
            ok = 0
        if not ok:
            raise HTTPException(status_code=400, detail="Incorrect password.")
    return reader


def _parse_page_sequence(value: str | None, total_pages: int) -> List[int]:
    if not value:
        return list(range(1, total_pages + 1))

    tokens = re.split(r"[,\s]+", value.strip())
    pages: List[int] = []
    for token in tokens:
        if not token:
            continue
        if "-" in token:
            start_str, end_str = token.split("-", 1)
            if not start_str or not end_str:
                raise HTTPException(status_code=400, detail="Invalid page range format.")
            try:
                start = int(start_str)
                end = int(end_str)
            except ValueError as exc:
                raise HTTPException(status_code=400, detail="Invalid page number.") from exc
            if start < 1 or end < 1 or start > total_pages or end > total_pages:
                raise HTTPException(status_code=400, detail="Page range out of bounds.")
            step = 1 if start <= end else -1
            pages.extend(list(range(start, end + step, step)))
        else:
            try:
                page = int(token)
            except ValueError as exc:
                raise HTTPException(status_code=400, detail="Invalid page number.") from exc
            if page < 1 or page > total_pages:
                raise HTTPException(status_code=400, detail="Page number out of bounds.")
            pages.append(page)

    if not pages:
        raise HTTPException(status_code=400, detail="No pages specified.")
    return pages


def _write_pdf_response(
    writer: PdfWriter,
    background_tasks: BackgroundTasks,
    name: str,
) -> FileResponse:
    output_path = OUTPUT_DIR / f"{name}_{uuid.uuid4().hex}.pdf"
    with output_path.open("wb") as f:
        writer.write(f)
    background_tasks.add_task(_safe_remove, output_path)
    return FileResponse(
        path=str(output_path),
        filename=f"{name}.pdf",
        media_type="application/pdf",
        background=background_tasks,
    )


def _clamp_pct(value: float, min_value: float = 0.0, max_value: float = 40.0) -> float:
    return max(min_value, min(max_value, value))


def _require_dependency(dep, name: str, hint: str) -> None:
    if dep is None:
        raise HTTPException(status_code=500, detail=f"{name} is required. {hint}")


def _html_to_pdf(html_path: Path) -> Path:
    _require_dependency(HTML, "WeasyPrint", "Install with: pip install weasyprint")
    output_path = OUTPUT_DIR / f"html_{uuid.uuid4().hex}.pdf"
    HTML(filename=str(html_path)).write_pdf(str(output_path))
    return output_path


def _render_text_pdf(title: str, paragraphs: List[str]) -> Path:
    _require_dependency(SimpleDocTemplate, "reportlab", "Install with: pip install reportlab")
    output_path = OUTPUT_DIR / f"report_{uuid.uuid4().hex}.pdf"
    doc = SimpleDocTemplate(str(output_path), pagesize=letter)
    styles = getSampleStyleSheet()
    story = [Paragraph(escape(title), styles["Heading1"]), Spacer(1, 12)]
    for para in paragraphs:
        text = (para or "").strip()
        if not text:
            story.append(Spacer(1, 8))
            continue
        safe_text = escape(text).replace("\n", "<br/>")
        story.append(Paragraph(safe_text, styles["BodyText"]))
        story.append(Spacer(1, 6))
    doc.build(story)
    return output_path


def _docx_to_pdf(source_path: Path) -> Path:
    _require_dependency(Document, "python-docx", "Install with: pip install python-docx")
    doc = Document(str(source_path))
    paragraphs: List[str] = []

    for para in doc.paragraphs:
        text = (para.text or "").strip()
        if text:
            paragraphs.append(text)

    for table in doc.tables:
        for row in table.rows:
            cells = [(cell.text or "").strip() for cell in row.cells]
            cleaned = [cell for cell in cells if cell]
            if cleaned:
                paragraphs.append(" | ".join(cleaned))

    return _render_text_pdf(source_path.stem, paragraphs or ["No readable text found in the DOCX file."])


def _pptx_to_pdf(source_path: Path) -> Path:
    _require_dependency(Presentation, "python-pptx", "Install with: pip install python-pptx")
    presentation = Presentation(str(source_path))
    paragraphs: List[str] = []

    for slide_number, slide in enumerate(presentation.slides, start=1):
        slide_texts: List[str] = []
        for shape in slide.shapes:
            text = getattr(shape, "text", "") or ""
            cleaned = text.strip()
            if cleaned:
                slide_texts.append(cleaned)
        paragraphs.append(f"Slide {slide_number}")
        paragraphs.extend(slide_texts or ["No readable text found on this slide."])
        paragraphs.append("")

    return _render_text_pdf(source_path.stem, paragraphs or ["No readable text found in the PPTX file."])


def _xlsx_to_pdf(source_path: Path) -> Path:
    _require_dependency(load_workbook, "openpyxl", "Install with: pip install openpyxl")
    workbook = load_workbook(str(source_path), data_only=True)
    paragraphs: List[str] = []

    try:
        for worksheet in workbook.worksheets:
            paragraphs.append(f"Sheet: {worksheet.title}")
            found_rows = False
            for row in worksheet.iter_rows(values_only=True):
                values = [str(cell).strip() for cell in row if cell is not None and str(cell).strip()]
                if values:
                    found_rows = True
                    paragraphs.append(" | ".join(values))
            if not found_rows:
                paragraphs.append("No readable data found on this sheet.")
            paragraphs.append("")
    finally:
        workbook.close()

    return _render_text_pdf(source_path.stem, paragraphs or ["No readable content found in the XLSX file."])


def _text_source_to_pdf(source_path: Path) -> Path:
    text = _read_text_source(source_path)
    return _render_text_pdf(source_path.stem, _split_text_blocks(text))


def _convert_source_to_pdf(source_path: Path) -> Path:
    ext = source_path.suffix.lower()
    if ext == ".pdf":
        return source_path
    if ext in {".png", ".jpg", ".jpeg", ".webp"}:
        return _image_to_pdf(source_path)
    if ext in {".html", ".htm"}:
        return _html_to_pdf(source_path)
    if ext == ".docx":
        return _docx_to_pdf(source_path)
    if ext == ".pptx":
        return _pptx_to_pdf(source_path)
    if ext == ".xlsx":
        return _xlsx_to_pdf(source_path)
    if ext in TEXT_SOURCE_EXTS:
        return _text_source_to_pdf(source_path)
    raise HTTPException(status_code=400, detail=f"Unsupported file type: {ext or 'unknown'}")


def _overlay_text_on_pages(
    reader: PdfReader,
    text_fn,
    position: str,
    font_size: int,
    angle: int = 0,
    opacity: float = 0.2,
    pages: Optional[List[int]] = None,
) -> PdfWriter:
    _require_dependency(rl_canvas, "reportlab", "Install with: pip install reportlab")
    writer = PdfWriter()
    target_pages = set(pages or range(1, len(reader.pages) + 1))
    margin = 36

    def draw_on_canvas(c, width, height, text):
        try:
            c.setFillAlpha(opacity)
        except Exception:
            pass
        c.setFont("Helvetica", font_size)

        if position == "top-left":
            x, y, align = margin, height - margin, "left"
        elif position == "top-center":
            x, y, align = width / 2, height - margin, "center"
        elif position == "top-right":
            x, y, align = width - margin, height - margin, "right"
        elif position == "bottom-left":
            x, y, align = margin, margin, "left"
        elif position == "bottom-center":
            x, y, align = width / 2, margin, "center"
        elif position == "bottom-right":
            x, y, align = width - margin, margin, "right"
        else:
            x, y, align = width / 2, height / 2, "center"

        c.saveState()
        if angle:
            c.translate(x, y)
            c.rotate(angle)
            x, y = 0, 0

        if align == "right":
            c.drawRightString(x, y, text)
        elif align == "center":
            c.drawCentredString(x, y, text)
        else:
            c.drawString(x, y, text)
        c.restoreState()

    for idx, page in enumerate(reader.pages, start=1):
        if idx in target_pages:
            width = float(page.mediabox.width)
            height = float(page.mediabox.height)
            packet = BytesIO()
            c = rl_canvas.Canvas(packet, pagesize=(width, height))
            draw_on_canvas(c, width, height, text_fn(idx))
            c.save()
            packet.seek(0)
            overlay = PdfReader(packet)
            page.merge_page(overlay.pages[0])
        writer.add_page(page)
    return writer


def _prepare_logo_bytes(logo_bytes: bytes, opacity: float) -> tuple[bytes, int, int]:
    with Image.open(BytesIO(logo_bytes)) as img:
        img = img.convert("RGBA")
        width, height = img.size
        opacity = max(0.05, min(1.0, float(opacity)))
        if opacity < 0.999:
            alpha = img.split()[3]
            alpha = alpha.point(lambda p: int(p * opacity))
            img.putalpha(alpha)
        output = BytesIO()
        img.save(output, format="PNG")
    return output.getvalue(), width, height


def _extract_text(reader: PdfReader) -> List[str]:
    texts: List[str] = []
    for page in reader.pages:
        try:
            text = page.extract_text() or ""
        except Exception:
            text = ""
        texts.append(text)
    return texts


def _normalize_text(text: str) -> str:
    return re.sub(r"\s+", " ", (text or "").strip().lower())


def _keyword_stats(text: str, top_n: int = 10) -> List[tuple[str, int]]:
    cleaned = re.sub(r"[^a-zA-Z0-9\s]", " ", text.lower())
    words = [w for w in cleaned.split() if len(w) > 3]
    stop = {
        "this", "that", "with", "from", "were", "have", "has", "will", "would", "could",
        "should", "about", "there", "their", "they", "them", "then", "than", "your", "yours",
        "into", "over", "also", "just", "more", "most", "some", "such", "when", "where", "what",
        "which", "while", "these", "those", "here", "there", "been", "being", "after", "before",
        "because", "between", "within", "without"
    }
    words = [w for w in words if w not in stop]
    counter = Counter(words)
    return counter.most_common(top_n)


def _chunk_text(text: str, size: int = 4000) -> List[str]:
    chunks: List[str] = []
    current = ""
    for line in text.splitlines():
        if len(current) + len(line) + 1 > size:
            if current:
                chunks.append(current)
            current = line
        else:
            current = f"{current}\n{line}" if current else line
    if current:
        chunks.append(current)
    return chunks


def _summarize_text(text: str, max_sentences: int = 5) -> List[str]:
    max_sentences = max(3, min(12, int(max_sentences)))
    sentences = [s.strip() for s in re.split(r"(?<=[.!?])\s+", text or "") if len(s.strip().split()) >= 5]
    if not sentences:
        return []

    keyword_weights = dict(_keyword_stats(text, top_n=30))
    scored: List[tuple[float, int, str]] = []
    for index, sentence in enumerate(sentences):
        tokens = re.findall(r"[a-zA-Z0-9']+", sentence.lower())
        score = sum(keyword_weights.get(token, 0) for token in tokens)
        score += min(len(tokens), 24) * 0.08
        scored.append((score, index, sentence))

    top_sentences = sorted(scored, key=lambda item: (-item[0], item[1]))[:max_sentences]
    return [sentence for _, _, sentence in sorted(top_sentences, key=lambda item: item[1])]


def _render_pdf_pages(
    pdf_bytes: bytes,
    dpi: int,
    first_page: int,
    last_page: int,
):
    if pdfium is not None:
        doc = pdfium.PdfDocument(pdf_bytes)
        try:
            pages = []
            scale = dpi / 72
            for page_index in range(first_page - 1, last_page):
                page = doc[page_index]
                bitmap = page.render(scale=scale)
                pages.append(bitmap.to_pil())
            return pages
        finally:
            close = getattr(doc, "close", None)
            if callable(close):
                close()

    try:
        return convert_from_bytes(
            pdf_bytes,
            dpi=dpi,
            first_page=first_page,
            last_page=last_page,
        )
    except (PDFInfoNotInstalledError, PDFPageCountError, PDFSyntaxError) as exc:
        raise HTTPException(
            status_code=500,
            detail="PDF render failed. Install pypdfium2 or Poppler.",
        ) from exc


@app.get("/")
def home():
    return FileResponse(str(INDEX_FILE))


@app.get("/api/health")
def health_check():
    return {"message": "PDF Tool API is running."}


@app.post("/merge")
async def merge_files(background_tasks: BackgroundTasks, files: List[UploadFile] = File(...)):
    if len(files) < 1:
        raise HTTPException(status_code=400, detail="Please upload at least 1 file.")

    temp_paths: List[Path] = []
    merge_inputs: List[Path] = []

    try:
        for upload in files:
            ext = Path(upload.filename or "").suffix.lower()
            if ext not in SUPPORTED_MERGE_EXTS:
                raise HTTPException(
                    status_code=400,
                    detail=(
                        f"Unsupported file: {upload.filename}. Allowed: PDF, images, HTML, DOCX, PPTX, XLSX, TXT, MD, CSV, JSON."
                    ),
                )

            saved_path = _save_upload(upload, UPLOAD_DIR)
            temp_paths.append(saved_path)
            pdf_path = _convert_source_to_pdf(saved_path)
            if pdf_path != saved_path:
                temp_paths.append(pdf_path)
            merge_inputs.append(pdf_path)

        merger = PdfMerger()
        for path in merge_inputs:
            merger.append(str(path))

        output_path = OUTPUT_DIR / f"merged_{uuid.uuid4().hex}.pdf"
        merger.write(str(output_path))
        merger.close()

        background_tasks.add_task(_safe_remove, output_path)

        return FileResponse(
            path=str(output_path),
            filename="merged.pdf",
            media_type="application/pdf",
            background=background_tasks,
        )
    finally:
        for path in temp_paths:
            _safe_remove(path)


@app.post("/convert-to-pdf")
async def convert_to_pdf(background_tasks: BackgroundTasks, files: List[UploadFile] = File(...)):
    return await merge_files(background_tasks=background_tasks, files=files)


@app.post("/pdf-info")
async def pdf_info(file: UploadFile = File(...)):
    if not (file.filename or "").lower().endswith(".pdf"):
        raise HTTPException(status_code=400, detail="Only PDF files are supported.")

    pdf_bytes = await file.read()
    if not pdf_bytes:
        raise HTTPException(status_code=400, detail="Uploaded PDF is empty.")

    try:
        pages = len(PdfReader(BytesIO(pdf_bytes)).pages)
    except Exception as exc:
        raise HTTPException(status_code=400, detail="Invalid PDF file.") from exc

    return {"pages": pages}


@app.post("/split")
async def split_part1(
    background_tasks: BackgroundTasks,
    file: UploadFile = File(...),
    page: int = Form(...),
):
    if not (file.filename or "").lower().endswith(".pdf"):
        raise HTTPException(status_code=400, detail="Only PDF files are supported for split.")

    pdf_bytes = await file.read()
    part1_path, part2_path, _ = _pdf_split_to_files(pdf_bytes, page)

    background_tasks.add_task(_safe_remove, part1_path)
    background_tasks.add_task(_safe_remove, part2_path)

    return FileResponse(
        path=str(part1_path),
        filename="part1.pdf",
        media_type="application/pdf",
        background=background_tasks,
    )


@app.post("/split/part2")
async def split_part2(
    background_tasks: BackgroundTasks,
    file: UploadFile = File(...),
    page: int = Form(...),
):
    if not (file.filename or "").lower().endswith(".pdf"):
        raise HTTPException(status_code=400, detail="Only PDF files are supported for split.")

    pdf_bytes = await file.read()
    part1_path, part2_path, _ = _pdf_split_to_files(pdf_bytes, page)

    background_tasks.add_task(_safe_remove, part1_path)
    background_tasks.add_task(_safe_remove, part2_path)

    return FileResponse(
        path=str(part2_path),
        filename="part2.pdf",
        media_type="application/pdf",
        background=background_tasks,
    )


@app.post("/pdf-to-image")
async def pdf_to_image(
    file: UploadFile = File(...),
    format: str = Form("png"),
    dpi: int = Form(300),
    page_from: Optional[int] = Form(None),
    page_to: Optional[int] = Form(None),
):
    if not (file.filename or "").lower().endswith(".pdf"):
        raise HTTPException(status_code=400, detail="Only PDF files are supported for conversion.")

    dpi = max(72, min(600, int(dpi)))
    requested_format, pil_format = _resolve_image_format(format)

    pdf_bytes = await file.read()
    if not pdf_bytes:
        raise HTTPException(status_code=400, detail="Uploaded PDF is empty.")

    try:
        total_pages = len(PdfReader(BytesIO(pdf_bytes)).pages)
    except Exception as exc:
        raise HTTPException(status_code=400, detail="Invalid PDF file.") from exc

    first_page = page_from if page_from is not None else 1
    last_page = page_to if page_to is not None else total_pages

    if first_page < 1 or first_page > total_pages:
        raise HTTPException(status_code=400, detail=f"page_from must be between 1 and {total_pages}.")
    if last_page < 1 or last_page > total_pages:
        raise HTTPException(status_code=400, detail=f"page_to must be between 1 and {total_pages}.")
    if first_page > last_page:
        raise HTTPException(status_code=400, detail="page_from cannot be greater than page_to.")

    images = _render_pdf_pages(
        pdf_bytes=pdf_bytes,
        dpi=dpi,
        first_page=first_page,
        last_page=last_page,
    )

    result = []
    for image in images:
        with BytesIO() as buf:
            export_image = image
            save_kwargs = {}

            if pil_format == "JPEG":
                if export_image.mode not in ("RGB", "L"):
                    export_image = export_image.convert("RGB")
                save_kwargs["quality"] = 92
            elif pil_format == "WEBP":
                if export_image.mode == "RGBA":
                    export_image = export_image.convert("RGB")
                save_kwargs["quality"] = 90

            export_image.save(buf, format=pil_format, **save_kwargs)
            result.append(base64.b64encode(buf.getvalue()).decode("ascii"))

    return {
        "format": requested_format,
        "page_from": first_page,
        "page_to": last_page,
        "images": result,
    }


@app.post("/remove-pages")
async def remove_pages(
    background_tasks: BackgroundTasks,
    file: UploadFile = File(...),
    pages: str = Form(...),
):
    if not (file.filename or "").lower().endswith(".pdf"):
        raise HTTPException(status_code=400, detail="Only PDF files are supported.")
    if not pages.strip():
        raise HTTPException(status_code=400, detail="Pages are required.")
    pdf_bytes = await file.read()
    reader = _read_pdf(pdf_bytes)
    total_pages = len(reader.pages)
    to_remove = set(_parse_page_sequence(pages, total_pages))

    writer = PdfWriter()
    for idx, page in enumerate(reader.pages, start=1):
        if idx not in to_remove:
            writer.add_page(page)

    if len(writer.pages) == 0:
        raise HTTPException(status_code=400, detail="All pages were removed.")
    return _write_pdf_response(writer, background_tasks, "removed_pages")


@app.post("/extract-pages")
async def extract_pages(
    background_tasks: BackgroundTasks,
    file: UploadFile = File(...),
    pages: str = Form(...),
):
    if not (file.filename or "").lower().endswith(".pdf"):
        raise HTTPException(status_code=400, detail="Only PDF files are supported.")
    if not pages.strip():
        raise HTTPException(status_code=400, detail="Pages are required.")
    pdf_bytes = await file.read()
    reader = _read_pdf(pdf_bytes)
    total_pages = len(reader.pages)
    order = _parse_page_sequence(pages, total_pages)

    writer = PdfWriter()
    for page_num in order:
        writer.add_page(reader.pages[page_num - 1])
    return _write_pdf_response(writer, background_tasks, "extracted_pages")


@app.post("/organize")
async def organize_pages(
    background_tasks: BackgroundTasks,
    file: UploadFile = File(...),
    pages: str = Form(...),
):
    if not (file.filename or "").lower().endswith(".pdf"):
        raise HTTPException(status_code=400, detail="Only PDF files are supported.")
    if not pages.strip():
        raise HTTPException(status_code=400, detail="Pages are required.")
    pdf_bytes = await file.read()
    reader = _read_pdf(pdf_bytes)
    total_pages = len(reader.pages)
    order = _parse_page_sequence(pages, total_pages)

    writer = PdfWriter()
    for page_num in order:
        writer.add_page(reader.pages[page_num - 1])
    return _write_pdf_response(writer, background_tasks, "organized")


@app.post("/rotate")
async def rotate_pages(
    background_tasks: BackgroundTasks,
    file: UploadFile = File(...),
    degrees: int = Form(90),
    pages: Optional[str] = Form(None),
):
    if not (file.filename or "").lower().endswith(".pdf"):
        raise HTTPException(status_code=400, detail="Only PDF files are supported.")
    pdf_bytes = await file.read()
    reader = _read_pdf(pdf_bytes)
    total_pages = len(reader.pages)

    if degrees % 90 != 0:
        raise HTTPException(status_code=400, detail="Rotation must be a multiple of 90.")
    target_pages = set(_parse_page_sequence(pages, total_pages)) if pages else set(range(1, total_pages + 1))

    writer = PdfWriter()
    for idx, page in enumerate(reader.pages, start=1):
        if idx in target_pages:
            if degrees >= 0:
                page.rotate_clockwise(degrees)
            else:
                page.rotate_counter_clockwise(abs(degrees))
        writer.add_page(page)
    return _write_pdf_response(writer, background_tasks, "rotated")


@app.post("/crop")
async def crop_pages(
    background_tasks: BackgroundTasks,
    file: UploadFile = File(...),
    pages: Optional[str] = Form(None),
    left: float = Form(0.0),
    right: float = Form(0.0),
    top: float = Form(0.0),
    bottom: float = Form(0.0),
):
    if not (file.filename or "").lower().endswith(".pdf"):
        raise HTTPException(status_code=400, detail="Only PDF files are supported.")
    pdf_bytes = await file.read()
    reader = _read_pdf(pdf_bytes)
    total_pages = len(reader.pages)

    left = _clamp_pct(left)
    right = _clamp_pct(right)
    top = _clamp_pct(top)
    bottom = _clamp_pct(bottom)
    if left + right >= 100 or top + bottom >= 100:
        raise HTTPException(status_code=400, detail="Crop percentages are too large.")

    target_pages = set(_parse_page_sequence(pages, total_pages)) if pages else set(range(1, total_pages + 1))
    writer = PdfWriter()
    for idx, page in enumerate(reader.pages, start=1):
        if idx in target_pages:
            box = page.mediabox
            width = float(box.width)
            height = float(box.height)
            llx = float(box.left) + width * (left / 100)
            lly = float(box.bottom) + height * (bottom / 100)
            urx = float(box.right) - width * (right / 100)
            ury = float(box.top) - height * (top / 100)
            if urx <= llx or ury <= lly:
                raise HTTPException(status_code=400, detail="Crop values remove the entire page.")
            page.cropbox.lower_left = (llx, lly)
            page.cropbox.upper_right = (urx, ury)
        writer.add_page(page)
    return _write_pdf_response(writer, background_tasks, "cropped")


@app.post("/compress")
async def compress_pdf(
    background_tasks: BackgroundTasks,
    file: UploadFile = File(...),
):
    if not (file.filename or "").lower().endswith(".pdf"):
        raise HTTPException(status_code=400, detail="Only PDF files are supported.")
    pdf_bytes = await file.read()
    reader = _read_pdf(pdf_bytes)
    writer = PdfWriter()
    for page in reader.pages:
        try:
            page.compress_content_streams()
        except Exception:
            pass
        writer.add_page(page)
    return _write_pdf_response(writer, background_tasks, "compressed")


@app.post("/optimize")
async def optimize_pdf(
    background_tasks: BackgroundTasks,
    file: UploadFile = File(...),
):
    if not (file.filename or "").lower().endswith(".pdf"):
        raise HTTPException(status_code=400, detail="Only PDF files are supported.")
    pdf_bytes = await file.read()
    reader = _read_pdf(pdf_bytes)
    writer = PdfWriter()
    for page in reader.pages:
        try:
            page.compress_content_streams()
        except Exception:
            pass
        writer.add_page(page)
    writer.add_metadata({})
    return _write_pdf_response(writer, background_tasks, "optimized")


@app.post("/repair")
async def repair_pdf(
    background_tasks: BackgroundTasks,
    file: UploadFile = File(...),
):
    if not (file.filename or "").lower().endswith(".pdf"):
        raise HTTPException(status_code=400, detail="Only PDF files are supported.")
    pdf_bytes = await file.read()
    reader = _read_pdf(pdf_bytes, strict=False)
    writer = PdfWriter()
    for page in reader.pages:
        writer.add_page(page)
    return _write_pdf_response(writer, background_tasks, "repaired")


@app.post("/protect")
async def protect_pdf(
    background_tasks: BackgroundTasks,
    file: UploadFile = File(...),
    password: str = Form(...),
):
    if not (file.filename or "").lower().endswith(".pdf"):
        raise HTTPException(status_code=400, detail="Only PDF files are supported.")
    if not password:
        raise HTTPException(status_code=400, detail="Password is required.")
    pdf_bytes = await file.read()
    if not pdf_bytes:
        raise HTTPException(status_code=400, detail="Uploaded PDF is empty.")
    try:
        reader = PdfReader(BytesIO(pdf_bytes))
    except Exception as exc:
        raise HTTPException(status_code=400, detail="Invalid PDF file.") from exc
    if reader.is_encrypted:
        raise HTTPException(status_code=400, detail="PDF is already protected. Unlock it first.")

    writer = PdfWriter()
    for page in reader.pages:
        writer.add_page(page)
    writer.encrypt(password)
    return _write_pdf_response(writer, background_tasks, "protected")


@app.post("/unlock")
async def unlock_pdf(
    background_tasks: BackgroundTasks,
    file: UploadFile = File(...),
    password: str = Form(...),
):
    if not (file.filename or "").lower().endswith(".pdf"):
        raise HTTPException(status_code=400, detail="Only PDF files are supported.")
    if not password:
        raise HTTPException(status_code=400, detail="Password is required.")
    pdf_bytes = await file.read()
    reader = _read_pdf(pdf_bytes, password=password)

    writer = PdfWriter()
    for page in reader.pages:
        writer.add_page(page)
    return _write_pdf_response(writer, background_tasks, "unlocked")


@app.post("/add-page-numbers")
async def add_page_numbers(
    background_tasks: BackgroundTasks,
    file: UploadFile = File(...),
    start: int = Form(1),
    position: str = Form("bottom-right"),
    pages: Optional[str] = Form(None),
):
    if not (file.filename or "").lower().endswith(".pdf"):
        raise HTTPException(status_code=400, detail="Only PDF files are supported.")
    pdf_bytes = await file.read()
    reader = _read_pdf(pdf_bytes)
    total_pages = len(reader.pages)
    target_pages = _parse_page_sequence(pages, total_pages) if pages else None
    start = max(1, int(start))

    writer = _overlay_text_on_pages(
        reader,
        lambda idx: str(start + idx - 1),
        position=position,
        font_size=10,
        angle=0,
        opacity=1.0,
        pages=target_pages,
    )
    return _write_pdf_response(writer, background_tasks, "page_numbers")


@app.post("/add-watermark")
async def add_watermark(
    background_tasks: BackgroundTasks,
    file: UploadFile = File(...),
    text: str = Form(...),
    opacity: float = Form(0.2),
    angle: int = Form(35),
    position: str = Form("center"),
    pages: Optional[str] = Form(None),
):
    if not (file.filename or "").lower().endswith(".pdf"):
        raise HTTPException(status_code=400, detail="Only PDF files are supported.")
    if not text.strip():
        raise HTTPException(status_code=400, detail="Watermark text is required.")
    pdf_bytes = await file.read()
    reader = _read_pdf(pdf_bytes)
    total_pages = len(reader.pages)
    target_pages = _parse_page_sequence(pages, total_pages) if pages else None

    base_font = 36
    if reader.pages:
        w = float(reader.pages[0].mediabox.width)
        h = float(reader.pages[0].mediabox.height)
        base_font = max(24, int(min(w, h) * 0.08))

    writer = _overlay_text_on_pages(
        reader,
        lambda idx: text,
        position=position,
        font_size=base_font,
        angle=int(angle),
        opacity=max(0.05, min(1.0, float(opacity))),
        pages=target_pages,
    )
    return _write_pdf_response(writer, background_tasks, "watermarked")


@app.post("/add-logo")
async def add_logo(
    background_tasks: BackgroundTasks,
    file: UploadFile = File(...),
    logo: UploadFile = File(...),
    scale: float = Form(0.2),
    opacity: float = Form(0.8),
    position: str = Form("bottom-right"),
    rotate: int = Form(0),
    pages: Optional[str] = Form(None),
):
    if not (file.filename or "").lower().endswith(".pdf"):
        raise HTTPException(status_code=400, detail="Only PDF files are supported.")
    if not (logo.filename or ""):
        raise HTTPException(status_code=400, detail="Logo image is required.")
    _require_dependency(fitz, "PyMuPDF", "Install with: pip install pymupdf")
    pdf_bytes = await file.read()
    logo_bytes_raw = await logo.read()
    if not logo_bytes_raw:
        raise HTTPException(status_code=400, detail="Logo image is required.")

    try:
        logo_bytes, logo_w, logo_h = _prepare_logo_bytes(logo_bytes_raw, opacity)
    except Exception:
        raise HTTPException(status_code=400, detail="Logo must be a valid image file.")

    doc = fitz.open(stream=pdf_bytes, filetype="pdf")
    total_pages = doc.page_count
    if total_pages == 0:
        doc.close()
        raise HTTPException(status_code=400, detail="PDF has no pages.")

    target_pages = _parse_page_sequence(pages, total_pages) if pages else list(range(1, total_pages + 1))
    scale = max(0.05, min(0.6, float(scale)))
    rotate = int(rotate)
    position = (position or "bottom-right").lower()
    valid_positions = {"top-left", "top-right", "bottom-left", "bottom-right", "center"}
    if position not in valid_positions:
        doc.close()
        raise HTTPException(status_code=400, detail="Invalid logo position.")

    for page_num in target_pages:
        page = doc[page_num - 1]
        rect = page.rect
        page_w, page_h = rect.width, rect.height
        margin = max(12, min(page_w, page_h) * 0.03)
        target_w = page_w * scale
        if logo_w > 0 and logo_h > 0:
            target_h = target_w * (logo_h / logo_w)
        else:
            target_h = target_w
        if target_h > page_h * 0.8:
            target_h = page_h * 0.8
            target_w = target_h * (logo_w / logo_h) if logo_h else target_h

        if position == "center":
            x = (page_w - target_w) / 2
            y = (page_h - target_h) / 2
        elif position == "top-left":
            x, y = margin, margin
        elif position == "top-right":
            x, y = page_w - target_w - margin, margin
        elif position == "bottom-left":
            x, y = margin, page_h - target_h - margin
        else:
            x, y = page_w - target_w - margin, page_h - target_h - margin

        logo_rect = fitz.Rect(x, y, x + target_w, y + target_h)
        page.insert_image(logo_rect, stream=logo_bytes, keep_proportion=True, overlay=True, rotate=rotate)

    output_path = OUTPUT_DIR / f"logo_{uuid.uuid4().hex}.pdf"
    doc.save(output_path)
    doc.close()
    background_tasks.add_task(_safe_remove, output_path)
    return FileResponse(
        path=str(output_path),
        filename="logo.pdf",
        media_type="application/pdf",
        background=background_tasks,
    )


@app.post("/sign")
async def sign_pdf(
    background_tasks: BackgroundTasks,
    file: UploadFile = File(...),
    text: str = Form(...),
):
    if not (file.filename or "").lower().endswith(".pdf"):
        raise HTTPException(status_code=400, detail="Only PDF files are supported.")
    if not text.strip():
        raise HTTPException(status_code=400, detail="Signature text is required.")
    pdf_bytes = await file.read()
    reader = _read_pdf(pdf_bytes)
    total_pages = len(reader.pages)
    if total_pages == 0:
        raise HTTPException(status_code=400, detail="PDF has no pages.")

    stamp = f"Signed by {text.strip()} on {datetime.now().strftime('%Y-%m-%d')}"
    writer = _overlay_text_on_pages(
        reader,
        lambda idx: stamp,
        position="bottom-right",
        font_size=11,
        angle=0,
        opacity=1.0,
        pages=[total_pages],
    )
    return _write_pdf_response(writer, background_tasks, "signed")


@app.post("/redact")
async def redact_pdf(
    background_tasks: BackgroundTasks,
    file: UploadFile = File(...),
    text: str = Form(...),
    pages: Optional[str] = Form(None),
):
    if not (file.filename or "").lower().endswith(".pdf"):
        raise HTTPException(status_code=400, detail="Only PDF files are supported.")
    if not text.strip():
        raise HTTPException(status_code=400, detail="Redaction text is required.")
    _require_dependency(fitz, "PyMuPDF", "Install with: pip install pymupdf")
    pdf_bytes = await file.read()
    doc = fitz.open(stream=pdf_bytes, filetype="pdf")
    total_pages = doc.page_count
    target_pages = _parse_page_sequence(pages, total_pages) if pages else list(range(1, total_pages + 1))

    for page_num in target_pages:
        page = doc[page_num - 1]
        for rect in page.search_for(text):
            page.add_redact_annot(rect, fill=(0, 0, 0))
        page.apply_redactions()

    output_path = OUTPUT_DIR / f"redacted_{uuid.uuid4().hex}.pdf"
    doc.save(output_path)
    doc.close()
    background_tasks.add_task(_safe_remove, output_path)
    return FileResponse(
        path=str(output_path),
        filename="redacted.pdf",
        media_type="application/pdf",
        background=background_tasks,
    )


@app.post("/pdf-to-word")
async def pdf_to_word(
    background_tasks: BackgroundTasks,
    file: UploadFile = File(...),
):
    if not (file.filename or "").lower().endswith(".pdf"):
        raise HTTPException(status_code=400, detail="Only PDF files are supported.")
    _require_dependency(Document, "python-docx", "Install with: pip install python-docx")
    pdf_bytes = await file.read()
    reader = _read_pdf(pdf_bytes)
    texts = _extract_text(reader)

    doc = Document()
    for idx, text in enumerate(texts, start=1):
        doc.add_heading(f"Page {idx}", level=2)
        for line in (text or "").splitlines():
            doc.add_paragraph(line)
    output_path = OUTPUT_DIR / f"pdf_to_word_{uuid.uuid4().hex}.docx"
    doc.save(output_path)
    background_tasks.add_task(_safe_remove, output_path)
    return FileResponse(
        path=str(output_path),
        filename="converted.docx",
        media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        background=background_tasks,
    )


@app.post("/pdf-to-excel")
async def pdf_to_excel(
    background_tasks: BackgroundTasks,
    file: UploadFile = File(...),
):
    if not (file.filename or "").lower().endswith(".pdf"):
        raise HTTPException(status_code=400, detail="Only PDF files are supported.")
    _require_dependency(Workbook, "openpyxl", "Install with: pip install openpyxl")
    pdf_bytes = await file.read()
    reader = _read_pdf(pdf_bytes)
    texts = _extract_text(reader)

    wb = Workbook()
    for idx, text in enumerate(texts, start=1):
        ws = wb.active if idx == 1 else wb.create_sheet()
        ws.title = f"Page {idx}"
        for line in (text or "").splitlines():
            ws.append([line])
    output_path = OUTPUT_DIR / f"pdf_to_excel_{uuid.uuid4().hex}.xlsx"
    wb.save(output_path)
    background_tasks.add_task(_safe_remove, output_path)
    return FileResponse(
        path=str(output_path),
        filename="converted.xlsx",
        media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        background=background_tasks,
    )


@app.post("/pdf-to-powerpoint")
async def pdf_to_powerpoint(
    background_tasks: BackgroundTasks,
    file: UploadFile = File(...),
):
    if not (file.filename or "").lower().endswith(".pdf"):
        raise HTTPException(status_code=400, detail="Only PDF files are supported.")
    _require_dependency(Presentation, "python-pptx", "Install with: pip install python-pptx")
    pdf_bytes = await file.read()
    reader = _read_pdf(pdf_bytes)
    total_pages = len(reader.pages)
    images = _render_pdf_pages(pdf_bytes, dpi=150, first_page=1, last_page=total_pages)

    prs = Presentation()
    blank_layout = prs.slide_layouts[6]
    for image in images:
        slide = prs.slides.add_slide(blank_layout)
        img_bytes = BytesIO()
        image.save(img_bytes, format="PNG")
        img_bytes.seek(0)
        slide.shapes.add_picture(img_bytes, 0, 0, width=prs.slide_width, height=prs.slide_height)

    output_path = OUTPUT_DIR / f"pdf_to_ppt_{uuid.uuid4().hex}.pptx"
    prs.save(output_path)
    background_tasks.add_task(_safe_remove, output_path)
    return FileResponse(
        path=str(output_path),
        filename="converted.pptx",
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        background=background_tasks,
    )


@app.post("/translate")
async def translate_pdf(
    background_tasks: BackgroundTasks,
    file: UploadFile = File(...),
    lang: str = Form(...),
):
    if not (file.filename or "").lower().endswith(".pdf"):
        raise HTTPException(status_code=400, detail="Only PDF files are supported.")
    _require_dependency(Translator, "googletrans", "Install with: pip install googletrans==4.0.0-rc1")
    pdf_bytes = await file.read()
    reader = _read_pdf(pdf_bytes)
    full_text = "\n\n".join(_extract_text(reader))

    translator = Translator()
    chunks = _chunk_text(full_text or "")
    translated_parts: List[str] = []
    for chunk in chunks:
        if not chunk.strip():
            continue
        try:
            translated_parts.append(translator.translate(chunk, dest=lang).text)
        except Exception as exc:
            raise HTTPException(
                status_code=500,
                detail="Translation failed. Check your internet connection and language code.",
            ) from exc

    translated_text = "\n\n".join(translated_parts)
    output_path = _render_text_pdf("Translated PDF", translated_text.split("\n\n"))
    background_tasks.add_task(_safe_remove, output_path)
    return FileResponse(
        path=str(output_path),
        filename="translated.pdf",
        media_type="application/pdf",
        background=background_tasks,
    )


@app.post("/compare")
async def compare_pdfs(
    background_tasks: BackgroundTasks,
    file: UploadFile = File(...),
    file2: UploadFile = File(...),
):
    if not (file.filename or "").lower().endswith(".pdf") or not (file2.filename or "").lower().endswith(".pdf"):
        raise HTTPException(status_code=400, detail="Only PDF files are supported.")
    pdf_bytes_1 = await file.read()
    pdf_bytes_2 = await file2.read()
    reader_a = _read_pdf(pdf_bytes_1)
    reader_b = _read_pdf(pdf_bytes_2)
    texts_a = _extract_text(reader_a)
    texts_b = _extract_text(reader_b)

    max_pages = max(len(texts_a), len(texts_b))
    diffs: List[int] = []
    for i in range(max_pages):
        text_a = texts_a[i] if i < len(texts_a) else ""
        text_b = texts_b[i] if i < len(texts_b) else ""
        if _normalize_text(text_a) != _normalize_text(text_b):
            diffs.append(i + 1)

    paragraphs = [
        f"PDF A pages: {len(texts_a)}",
        f"PDF B pages: {len(texts_b)}",
        "Different pages: " + (", ".join(map(str, diffs)) if diffs else "None"),
    ]
    output_path = _render_text_pdf("PDF Compare Report", paragraphs)
    background_tasks.add_task(_safe_remove, output_path)
    return FileResponse(
        path=str(output_path),
        filename="compare_report.pdf",
        media_type="application/pdf",
        background=background_tasks,
    )


@app.post("/pdf-intelligence")
async def pdf_intelligence(
    background_tasks: BackgroundTasks,
    file: UploadFile = File(...),
):
    if not (file.filename or "").lower().endswith(".pdf"):
        raise HTTPException(status_code=400, detail="Only PDF files are supported.")
    pdf_bytes = await file.read()
    reader = _read_pdf(pdf_bytes)
    texts = _extract_text(reader)
    full_text = " ".join(texts)
    words = full_text.split()
    keywords = _keyword_stats(full_text)

    paragraphs = [
        f"Total pages: {len(reader.pages)}",
        f"Total words: {len(words)}",
        f"Total characters: {len(full_text)}",
        "Top keywords: " + ", ".join([f"{k} ({v})" for k, v in keywords]) if keywords else "Top keywords: None",
    ]
    output_path = _render_text_pdf("PDF Intelligence Report", paragraphs)
    background_tasks.add_task(_safe_remove, output_path)
    return FileResponse(
        path=str(output_path),
        filename="pdf_intelligence.pdf",
        media_type="application/pdf",
        background=background_tasks,
    )


@app.post("/summarize")
async def summarize_pdf(
    background_tasks: BackgroundTasks,
    file: UploadFile = File(...),
    sentences: int = Form(5),
):
    if not (file.filename or "").lower().endswith(".pdf"):
        raise HTTPException(status_code=400, detail="Only PDF files are supported.")

    pdf_bytes = await file.read()
    reader = _read_pdf(pdf_bytes)
    full_text = "\n\n".join(_extract_text(reader)).strip()
    if not full_text:
        raise HTTPException(status_code=400, detail="No readable text found in the PDF.")

    summary_lines = _summarize_text(full_text, max_sentences=sentences)
    keywords = _keyword_stats(full_text, top_n=8)
    paragraphs = [
        f"Pages: {len(reader.pages)}",
        f"Words: {len(full_text.split())}",
        "",
        "Summary",
        *(summary_lines or ["No strong summary sentences were detected."]),
        "",
        "Top keywords: " + ", ".join(f"{word} ({count})" for word, count in keywords) if keywords else "Top keywords: None",
    ]

    output_path = _render_text_pdf("PDF Summary Report", paragraphs)
    background_tasks.add_task(_safe_remove, output_path)
    return FileResponse(
        path=str(output_path),
        filename="summary.pdf",
        media_type="application/pdf",
        background=background_tasks,
    )


@app.post("/scan-to-pdf")
async def scan_to_pdf(
    background_tasks: BackgroundTasks,
    files: List[UploadFile] = File(...),
):
    if not files:
        raise HTTPException(status_code=400, detail="Please upload at least 1 image.")

    images: List[Image.Image] = []
    for upload in files:
        ext = Path(upload.filename or "").suffix.lower()
        if ext not in {".png", ".jpg", ".jpeg", ".webp"}:
            raise HTTPException(status_code=400, detail="Only image files are supported for scanning.")
        img_bytes = await upload.read()
        with Image.open(BytesIO(img_bytes)) as img:
            if img.mode in ("RGBA", "LA"):
                bg = Image.new("RGB", img.size, "white")
                bg.paste(img, mask=img.getchannel("A"))
                images.append(bg)
            else:
                images.append(img.convert("RGB"))

    if not images:
        raise HTTPException(status_code=400, detail="No valid images found.")

    output_path = OUTPUT_DIR / f"scan_{uuid.uuid4().hex}.pdf"
    first, rest = images[0], images[1:]
    first.save(output_path, "PDF", save_all=True, append_images=rest)
    background_tasks.add_task(_safe_remove, output_path)
    return FileResponse(
        path=str(output_path),
        filename="scan.pdf",
        media_type="application/pdf",
        background=background_tasks,
    )


if __name__ == "__main__":
    host = os.getenv("PDF_TOOL_HOST", "0.0.0.0")
    port = int(os.getenv("PORT", os.getenv("PDF_TOOL_PORT", "8001")))
    open_browser = _env_flag("PDF_TOOL_OPEN_BROWSER", default=not _is_hosted_environment())
    browser_host = "localhost" if host == "0.0.0.0" else host
    app_url = f"http://{browser_host}:{port}"

    if open_browser:
        try:
            webbrowser.open_new_tab(app_url)
        except Exception:
            pass

    uvicorn.run(
        app,
        host=host,
        port=port,
        reload=False,
    )
